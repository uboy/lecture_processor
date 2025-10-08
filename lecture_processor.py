#!/usr/bin/env python3
"""
lecture_processor.py

Принимает видео/аудио файл и (опционально) слайды лекции (pdf/pptx).

Выполняет:
 - транскрипцию речи с помощью Whisper (локально);
 - извлечение из слайдов текста и добавление его в общий контекст;
 - разделение лекции на основную часть и Q&A;
 - создание конспекта с ключевыми понятиями, определениями, списком литературы,
   а также вопросами для самопроверки и ответами.

Выходные файлы:
 - full_transcript.txt
 - lecture_only.txt
 - q_and_a.txt
 - notes.txt

Пример запуска:
 python3 lecture_processor.py --input lecture.mp4 --slides slides.pdf --whisper-model small --ollama-model llama2 --outdir ./out

Зависимости:
 pip install -U openai-whisper ffmpeg-python requests tqdm python-pptx PyMuPDF
 sudo apt install ffmpeg

"""

import argparse
import subprocess
import os
import json
import tempfile
from pathlib import Path
from typing import Tuple

import requests
from tqdm import tqdm

try:
    import whisper
except Exception:
    whisper = None

# PDF и PPTX парсеры
import fitz  # PyMuPDF
from pptx import Presentation


def extract_audio(input_path: str, out_audio: str) -> None:
    cmd = ["ffmpeg", "-y", "-i", input_path, "-ar", "16000", "-ac", "1", "-vn", out_audio]
    subprocess.check_call(cmd)


def transcribe_with_whisper(audio_path: str, model_name: str = "large-v3") -> dict:
    if whisper is None:
        raise RuntimeError("whisper python package not available. Install openai-whisper")
    print(f"Loading Whisper model '{model_name}'...")
    model = whisper.load_model(model_name)
    print("Transcribing...")
    result = model.transcribe(audio_path, verbose=False)
    return result


def extract_text_from_slides(slide_path: str) -> str:
    path = Path(slide_path)
    text = ""
    if path.suffix.lower() == ".pdf":
        with fitz.open(slide_path) as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    elif path.suffix.lower() == ".pptx":
        prs = Presentation(slide_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        print(f"Unsupported slide format: {path.suffix}")
    return text.strip()


def save_text(path: str, text: str) -> None:
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)


def load_text(path: str) -> str:
    """Load text from file if it exists."""
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def call_ollama(prompt: str, ollama_model: str = "llama2", temperature: float = 0.0, max_tokens: int = 2000) -> str:
    """Call Ollama API with proper handling of streaming responses."""
    url = "http://localhost:11434/api/generate"
    headers = {"Content-Type": "application/json"}
    payload = {
        "model": ollama_model,
        "prompt": prompt,
        "stream": True,  # Use streaming for reliability
        "options": {
            "num_predict": max_tokens,
            "temperature": temperature
        }
    }

    try:
        # Use stream=True to get chunks as they arrive
        resp = requests.post(url, headers=headers, json=payload, stream=True, timeout=600)
        resp.raise_for_status()

        full_response = ""
        print("Receiving response from Ollama...", end="", flush=True)

        # Process streaming response line by line
        for line in resp.iter_lines():
            if line:
                try:
                    chunk = json.loads(line.decode('utf-8'))

                    # Extract response text from chunk
                    if "response" in chunk:
                        full_response += chunk["response"]
                        print(".", end="", flush=True)  # Progress indicator

                    # Check if this is the final chunk
                    if chunk.get("done", False):
                        print()  # New line after completion
                        break

                except json.JSONDecodeError as e:
                    print(f"\nWarning: Failed to parse chunk: {line[:100]}")
                    continue

        if not full_response:
            raise RuntimeError("Empty response from Ollama")

        return full_response.strip()

    except requests.exceptions.Timeout:
        raise RuntimeError("Ollama request timed out after 600 seconds")
    except requests.exceptions.RequestException as e:
        raise RuntimeError(f"Failed to call Ollama API: {e}")
    except Exception as e:
        raise RuntimeError(f"Error processing Ollama response: {e}")


def extract_json_from_response(text) -> dict:
    """Extract JSON from response that may contain thinking tags or markdown."""
    import re

    # If already a dict, return it
    if isinstance(text, dict):
        return text

    # Ensure text is a string
    if not isinstance(text, str):
        text = str(text)

    # Remove <think> tags and their content
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)

    # Remove markdown code blocks
    text = re.sub(r'```json\s*', '', text)
    text = re.sub(r'```\s*', '', text)

    # Try to find JSON object in the text
    # Look for content between { and }
    start = text.find('{')
    if start == -1:
        raise ValueError("No JSON object found in response")

    # Find matching closing brace
    brace_count = 0
    end = start
    for i in range(start, len(text)):
        if text[i] == '{':
            brace_count += 1
        elif text[i] == '}':
            brace_count -= 1
            if brace_count == 0:
                end = i + 1
                break

    if end > start:
        json_str = text[start:end]
        return json.loads(json_str)

    raise ValueError("Could not extract valid JSON from response")


def split_transcript_with_ollama(full_text: str, slide_text: str, ollama_model: str) -> Tuple[str, str]:
    """Split transcript into lecture and Q&A sections."""
    # Truncate inputs if too long to avoid context length issues
    max_transcript_chars = 15000
    max_slides_chars = 5000

    truncated_transcript = full_text[:max_transcript_chars]
    if len(full_text) > max_transcript_chars:
        truncated_transcript += "\n\n[... transcript truncated ...]"

    truncated_slides = slide_text[:max_slides_chars]
    if len(slide_text) > max_slides_chars:
        truncated_slides += "\n\n[... slides truncated ...]"

    prompt = (
        "Split the lecture transcript into two parts: main lecture and Q&A section.\n\n"
        "IMPORTANT: Output ONLY a JSON object, nothing else. No explanations, no thinking, just JSON.\n\n"
        "Required JSON format:\n"
        '{"lecture": "main lecture content here", "q_and_a": "Q&A section here"}\n\n'
        "If there is no Q&A section, use empty string for q_and_a.\n\n"
        f"Slides text (reference):\n{truncated_slides}\n\n"
        f"Transcript to split:\n{truncated_transcript}\n\n"
        "JSON output:"
    )

    print("Calling Ollama to split transcript...")
    out = call_ollama(prompt, ollama_model=ollama_model, max_tokens=8000, temperature=0.1)

    try:
        # Extract JSON from response (handles <think> tags, markdown, etc)
        json_str = extract_json_from_response(out)

        obj = json.loads(json_str)
        lecture = obj.get("lecture", "")
        q_and_a = obj.get("q_and_a", "")

        if isinstance(q_and_a, list):
            q_and_a_text = "\n\n".join([f"Q: {i.get('question', '')}\nA: {i.get('answer', '')}" for i in q_and_a])
        else:
            q_and_a_text = str(q_and_a)

        return lecture.strip(), q_and_a_text.strip()
    except Exception as e:
        print(f"Warning: Failed to parse JSON response: {e}")
        print(f"Raw response (first 1000 chars): {out[:1000]}")
        # Fallback: return full text as lecture, empty Q&A
        return full_text.strip(), ""


def generate_notes_with_ollama(lecture_text: str, q_and_a_text: str, slide_text: str, ollama_model: str) -> str:
    """Generate study notes from lecture content."""
    # Truncate inputs to manageable size
    max_lecture_chars = 10000
    max_qa_chars = 3000
    max_slides_chars = 5000

    truncated_lecture = lecture_text[:max_lecture_chars]
    if len(lecture_text) > max_lecture_chars:
        truncated_lecture += "\n\n[... lecture truncated ...]"

    truncated_qa = q_and_a_text[:max_qa_chars]
    if len(q_and_a_text) > max_qa_chars:
        truncated_qa += "\n\n[... Q&A truncated ...]"

    truncated_slides = slide_text[:max_slides_chars]
    if len(slide_text) > max_slides_chars:
        truncated_slides += "\n\n[... slides truncated ...]"

    prompt = (
        "Create detailed study notes for this lecture.\n\n"
        "IMPORTANT: Output your notes directly. Do NOT use <think> tags or JSON format.\n\n"
        "Include these sections:\n"
        "1) Summary (3-6 paragraphs)\n"
        "2) Key concepts and definitions\n"
        "3) Suggested reading list (5-10 items)\n"
        "4) 10 practice questions\n"
        "5) Answers to the questions\n\n"
        f"Lecture:\n{truncated_lecture}\n\n"
        f"Q&A:\n{truncated_qa}\n\n"
        f"Slides:\n{truncated_slides}\n\n"
        "Study notes:"
    )

    print("Calling Ollama to generate notes...")
    response = call_ollama(prompt, ollama_model=ollama_model, max_tokens=8000, temperature=0.3)

    # Remove <think> tags if present
    import re
    response = re.sub(r'<think>.*?</think>', '', response, flags=re.DOTALL)

    return response.strip()


def main():
    parser = argparse.ArgumentParser(
        description="Process lecture (audio/video + optional slides) into transcript, Q&A, and notes.")
    parser.add_argument("--input", "-i", required=True, help="Input audio or video file")
    parser.add_argument("--slides", "-s", help="Optional slide file (.pdf or .pptx)")
    parser.add_argument("--whisper-model", default="large-v3", help="Whisper model size")
    parser.add_argument("--ollama-model", default="llama2", help="Ollama model name")
    parser.add_argument("--outdir", default="./output", help="Output directory")
    parser.add_argument("--no-audio-extract", action="store_true", help="Skip audio extraction if already audio")
    parser.add_argument("--force-transcribe", action="store_true",
                        help="Force re-transcription even if transcript exists")
    args = parser.parse_args()

    Path(args.outdir).mkdir(parents=True, exist_ok=True)

    # Check if transcript already exists
    transcript_path = os.path.join(args.outdir, "full_transcript.txt")
    full_text = ""

    if os.path.exists(transcript_path) and not args.force_transcribe:
        print(f"Loading existing transcript from {transcript_path}")
        full_text = load_text(transcript_path)

    if not full_text:
        audio_path = os.path.join(tempfile.gettempdir(), "lecture_audio.wav")
        if not args.no_audio_extract:
            print("Extracting audio...")
            extract_audio(args.input, audio_path)
        else:
            audio_path = args.input

        print("Transcribing with Whisper...")
        result = transcribe_with_whisper(audio_path, model_name=args.whisper_model)
        full_text = result.get("text", "")
        save_text(transcript_path, full_text)
        print(f"Transcript saved to {transcript_path}")

    # Extract slides if provided
    slide_text = ""
    if args.slides:
        slides_text_path = os.path.join(args.outdir, "slides_text.txt")
        if os.path.exists(slides_text_path):
            print(f"Loading existing slides text from {slides_text_path}")
            slide_text = load_text(slides_text_path)
        else:
            print(f"Extracting text from slides: {args.slides}")
            slide_text = extract_text_from_slides(args.slides)
            save_text(slides_text_path, slide_text)

    # Split into lecture and Q&A
    print("Splitting into lecture/Q&A...")
    lecture_text, q_and_a_text = split_transcript_with_ollama(full_text, slide_text, args.ollama_model)
    save_text(os.path.join(args.outdir, "lecture_only.txt"), lecture_text)
    save_text(os.path.join(args.outdir, "q_and_a.txt"), q_and_a_text)

    # Generate notes
    print("Generating notes...")
    notes = generate_notes_with_ollama(lecture_text, q_and_a_text, slide_text, args.ollama_model)
    save_text(os.path.join(args.outdir, "notes.txt"), notes)

    print(f"Done! Results in {args.outdir}")


if __name__ == "__main__":
    main()